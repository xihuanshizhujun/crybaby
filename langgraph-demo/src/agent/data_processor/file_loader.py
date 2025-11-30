"""文件加载器 - 基于 deep-searcher 实现"""
import os
import logging
from pathlib import Path
from typing import List, Dict, Any, Optional
import mimetypes

logger = logging.getLogger(__name__)


class FileLoader:
    """文件加载器，支持多种文件格式"""
    
    SUPPORTED_EXTENSIONS = {
        '.pdf', '.docx', '.doc', '.pptx', '.ppt',
        '.xlsx', '.xls', '.txt', '.md', '.csv'
    }
    
    def __init__(self):
        self.loaders = {}
        self._init_loaders()
    
    def _init_loaders(self):
        """初始化各种文件加载器"""
        # PDF
        try:
            import pdfplumber
            self.loaders['.pdf'] = self._load_pdf
        except ImportError:
            try:
                import pypdf
                self.loaders['.pdf'] = self._load_pdf_pypdf
            except ImportError:
                logger.warning("PDF loader not available")
        
        # DOCX
        try:
            from docx import Document
            self.loaders['.docx'] = self._load_docx
        except ImportError:
            logger.warning("DOCX loader not available")
        
        # DOC
        try:
            from docx2python import docx2python
            self.loaders['.doc'] = self._load_doc
        except ImportError:
            logger.warning("DOC loader not available")
        
        # PPTX
        try:
            from pptx import Presentation
            self.loaders['.pptx'] = self._load_pptx
            self.loaders['.ppt'] = self._load_pptx
        except ImportError:
            logger.warning("PPTX loader not available")
        
        # Excel
        try:
            import pandas as pd
            self.loaders['.xlsx'] = self._load_excel
            self.loaders['.xls'] = self._load_excel
            self.loaders['.csv'] = self._load_csv
        except ImportError:
            logger.warning("Excel loader not available")
        
        # Text
        self.loaders['.txt'] = self._load_text
        self.loaders['.md'] = self._load_text
    
    def is_supported(self, file_path: str) -> bool:
        """检查文件是否支持"""
        ext = Path(file_path).suffix.lower()
        return ext in self.SUPPORTED_EXTENSIONS and ext in self.loaders
    
    def load(self, file_path: str) -> Dict[str, Any]:
        """加载文件
        
        Args:
            file_path: 文件路径
            
        Returns:
            包含 text 和 metadata 的字典
        """
        if not os.path.exists(file_path):
            raise FileNotFoundError(f"File not found: {file_path}")
        
        ext = Path(file_path).suffix.lower()
        if ext not in self.loaders:
            raise ValueError(f"Unsupported file type: {ext}")
        
        loader = self.loaders[ext]
        result = loader(file_path)
        
        # 添加元数据
        result['metadata'] = {
            'file_path': file_path,
            'file_name': Path(file_path).name,
            'file_type': ext,
            'file_size': os.path.getsize(file_path),
        }
        
        return result
    
    def _load_pdf(self, file_path: str) -> Dict[str, Any]:
        """加载PDF文件（使用pdfplumber）"""
        import pdfplumber
        
        text_parts = []
        with pdfplumber.open(file_path) as pdf:
            for page in pdf.pages:
                text = page.extract_text()
                if text:
                    text_parts.append(text)
        
        return {
            'text': '\n\n'.join(text_parts),
            'page_count': len(text_parts),
        }
    
    def _load_pdf_pypdf(self, file_path: str) -> Dict[str, Any]:
        """加载PDF文件（使用pypdf，备选方案）"""
        import pypdf
        
        text_parts = []
        with open(file_path, 'rb') as f:
            reader = pypdf.PdfReader(f)
            for page in reader.pages:
                text = page.extract_text()
                if text:
                    text_parts.append(text)
        
        return {
            'text': '\n\n'.join(text_parts),
            'page_count': len(text_parts),
        }
    
    def _load_docx(self, file_path: str) -> Dict[str, Any]:
        """加载DOCX文件"""
        from docx import Document
        
        doc = Document(file_path)
        text_parts = []
        
        for paragraph in doc.paragraphs:
            if paragraph.text.strip():
                text_parts.append(paragraph.text)
        
        return {
            'text': '\n\n'.join(text_parts),
        }
    
    def _load_doc(self, file_path: str) -> Dict[str, Any]:
        """加载DOC文件"""
        from docx2python import docx2python
        
        doc = docx2python(file_path)
        text = doc.text
        
        return {
            'text': text,
        }
    
    def _load_pptx(self, file_path: str) -> Dict[str, Any]:
        """加载PPTX文件"""
        from pptx import Presentation
        
        prs = Presentation(file_path)
        text_parts = []
        
        for slide_num, slide in enumerate(prs.slides):
            slide_text = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_text.append(shape.text)
            
            if slide_text:
                text_parts.append(f"\n幻灯片 {slide_num + 1}:\n" + "\n".join(slide_text))
        
        return {
            'text': '\n\n'.join(text_parts),
            'slide_count': len(prs.slides),
        }
    
    def _load_excel(self, file_path: str) -> Dict[str, Any]:
        """加载Excel文件"""
        import pandas as pd
        
        text_parts = []
        excel_file = pd.ExcelFile(file_path)
        
        for sheet_name in excel_file.sheet_names:
            df = pd.read_excel(excel_file, sheet_name=sheet_name)
            df = df.fillna('')
            
            text_parts.append(f"\n工作表: {sheet_name}\n")
            text_parts.append(df.to_string(index=False))
        
        return {
            'text': '\n\n'.join(text_parts),
            'sheet_count': len(excel_file.sheet_names),
        }
    
    def _load_csv(self, file_path: str) -> Dict[str, Any]:
        """加载CSV文件"""
        import pandas as pd
        
        df = pd.read_csv(file_path)
        df = df.fillna('')
        
        return {
            'text': df.to_string(index=False),
        }
    
    def _load_text(self, file_path: str) -> Dict[str, Any]:
        """加载文本文件"""
        encodings = ['utf-8', 'gbk', 'gb2312', 'latin-1']
        
        for encoding in encodings:
            try:
                with open(file_path, 'r', encoding=encoding) as f:
                    text = f.read()
                return {
                    'text': text,
                    'encoding': encoding,
                }
            except UnicodeDecodeError:
                continue
        
        # 如果所有编码都失败，使用错误处理
        with open(file_path, 'r', encoding='utf-8', errors='replace') as f:
            text = f.read()
        return {
            'text': text,
            'encoding': 'utf-8 (with errors replaced)',
        }


# 全局实例
_file_loader = None

def get_file_loader() -> FileLoader:
    """获取文件加载器单例"""
    global _file_loader
    if _file_loader is None:
        _file_loader = FileLoader()
    return _file_loader

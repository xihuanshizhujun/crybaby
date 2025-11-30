"""统一文件管理器，支持多种文件格式的解析和处理"""

import logging
import chardet
from typing import Dict, Any, Optional
from pathlib import Path
import pandas as pd

logger = logging.getLogger(__name__)


class FileManager:
    """统一文件管理器
    
    支持的文件格式：
    - PDF: 使用pdfplumber（更好的中文和表格支持）
    - DOCX/DOC: 使用python-docx和docx2python
    - PPTX/PPT: 使用python-pptx
    - Excel: 使用pandas/openpyxl
    - TXT: 自动检测编码
    """
    
    # 支持的文件扩展名
    SUPPORTED_EXTENSIONS = {
        ".pdf": "pdf",
        ".docx": "docx",
        ".doc": "doc",
        ".pptx": "pptx",
        ".ppt": "ppt",
        ".xlsx": "excel",
        ".xls": "excel",
        ".txt": "txt",
    }
    
    @staticmethod
    def is_supported(file_path: str) -> bool:
        """检查文件格式是否支持
        
        Args:
            file_path: 文件路径
            
        Returns:
            是否支持
        """
        ext = Path(file_path).suffix.lower()
        return ext in FileManager.SUPPORTED_EXTENSIONS
    
    @staticmethod
    def get_file_type(file_path: str) -> Optional[str]:
        """获取文件类型
        
        Args:
            file_path: 文件路径
            
        Returns:
            文件类型，如果不支持则返回None
        """
        ext = Path(file_path).suffix.lower()
        return FileManager.SUPPORTED_EXTENSIONS.get(ext)
    
    @staticmethod
    def parse_file(file_path: str) -> Dict[str, Any]:
        """解析文件，自动根据扩展名选择解析器
        
        Args:
            file_path: 文件路径
            
        Returns:
            解析后的文档内容，包含text、tables、metadata等
            
        Raises:
            ValueError: 不支持的文件格式
            Exception: 解析失败
        """
        if not FileManager.is_supported(file_path):
            raise ValueError(f"不支持的文件格式: {Path(file_path).suffix}")
        
        file_type = FileManager.get_file_type(file_path)
        
        try:
            if file_type == "pdf":
                return FileManager._parse_pdf(file_path)
            elif file_type == "docx":
                return FileManager._parse_docx(file_path)
            elif file_type == "doc":
                return FileManager._parse_doc(file_path)
            elif file_type in ["pptx", "ppt"]:
                return FileManager._parse_pptx(file_path)
            elif file_type == "excel":
                return FileManager._parse_excel(file_path)
            elif file_type == "txt":
                return FileManager._parse_txt(file_path)
            else:
                raise ValueError(f"未实现的文件类型: {file_type}")
        except Exception as e:
            logger.error(f"解析文件失败 {file_path}: {e}", exc_info=True)
            raise
    
    @staticmethod
    def _parse_pdf(file_path: str) -> Dict[str, Any]:
        """解析PDF文件（使用pdfplumber，更好的中文支持）
        
        Args:
            file_path: PDF文件路径
            
        Returns:
            解析后的内容
        """
        try:
            import pdfplumber
            
            text_content = []
            tables = []
            
            with pdfplumber.open(file_path) as pdf:
                for page_num, page in enumerate(pdf.pages):
                    # 提取文本（pdfplumber对中文支持更好）
                    page_text = page.extract_text()
                    if page_text:
                        text_content.append(page_text)
                    
                    # 提取表格（pdfplumber的表格提取能力很强）
                    page_tables = page.extract_tables()
                    for table_idx, table in enumerate(page_tables):
                        if table:
                            # 转换为DataFrame格式
                            try:
                                # 第一行作为表头
                                if len(table) > 1:
                                    df = pd.DataFrame(table[1:], columns=table[0])
                                    tables.append({
                                        "data": df.to_dict(orient='records'),
                                        "columns": df.columns.tolist(),
                                        "page": page_num + 1,
                                        "table_index": table_idx,
                                    })
                                else:
                                    # 只有一行，直接保存
                                    tables.append({
                                        "data": table,
                                        "columns": table[0] if table else [],
                                        "page": page_num + 1,
                                        "table_index": table_idx,
                                    })
                            except Exception as e:
                                logger.warning(f"转换表格失败 (页{page_num+1}, 表{table_idx}): {e}")
                                # 保留原始表格数据
                                tables.append({
                                    "data": table,
                                    "page": page_num + 1,
                                    "table_index": table_idx,
                                })
            
            return {
                "text": "\n\n".join(text_content),
                "tables": tables,
                "page_count": len(text_content),
                "file_type": "pdf",
            }
        except ImportError:
            logger.warning("pdfplumber未安装，尝试使用pypdf作为备选")
            # 备选方案：使用pypdf
            return FileManager._parse_pdf_fallback(file_path)
    
    @staticmethod
    def _parse_pdf_fallback(file_path: str) -> Dict[str, Any]:
        """PDF解析备选方案（使用pypdf）
        
        Args:
            file_path: PDF文件路径
            
        Returns:
            解析后的内容
        """
        import pypdf
        
        text_content = []
        tables = []
        
        with open(file_path, 'rb') as file:
            pdf_reader = pypdf.PdfReader(file)
            
            for page_num, page in enumerate(pdf_reader.pages):
                page_text = page.extract_text()
                if page_text:
                    text_content.append(page_text)
        
        return {
            "text": "\n\n".join(text_content),
            "tables": tables,
            "page_count": len(text_content),
            "file_type": "pdf",
        }
    
    @staticmethod
    def _parse_docx(file_path: str) -> Dict[str, Any]:
        """解析DOCX文件
        
        Args:
            file_path: DOCX文件路径
            
        Returns:
            解析后的内容
        """
        from docx import Document
        
        doc = Document(file_path)
        text_content = []
        tables = []
        
        # 提取文本
        for paragraph in doc.paragraphs:
            if paragraph.text.strip():
                text_content.append(paragraph.text)
        
        # 提取表格
        for table_idx, table in enumerate(doc.tables):
            table_data = []
            for row in table.rows:
                row_data = [cell.text.strip() for cell in row.cells]
                table_data.append(row_data)
            
            if table_data:
                try:
                    # 第一行作为表头
                    if len(table_data) > 1:
                        df = pd.DataFrame(table_data[1:], columns=table_data[0])
                        tables.append({
                            "data": df.to_dict(orient='records'),
                            "columns": df.columns.tolist(),
                            "table_index": table_idx,
                        })
                    else:
                        tables.append({
                            "data": table_data,
                            "columns": table_data[0] if table_data else [],
                            "table_index": table_idx,
                        })
                except Exception as e:
                    logger.warning(f"转换表格失败 (表{table_idx}): {e}")
                    tables.append({"data": table_data, "table_index": table_idx})
        
        return {
            "text": "\n\n".join(text_content),
            "tables": tables,
            "file_type": "docx",
        }
    
    @staticmethod
    def _parse_doc(file_path: str) -> Dict[str, Any]:
        """解析DOC文件（老格式Word文档）
        
        Args:
            file_path: DOC文件路径
            
        Returns:
            解析后的内容
        """
        try:
            # 尝试使用docx2python
            from docx2python import docx2python
            
            doc = docx2python(file_path)
            text_content = []
            tables = []
            
            # 提取文本
            text_content.append(doc.text)
            
            # 提取表格
            for table_idx, table in enumerate(doc.tables):
                if table:
                    try:
                        # docx2python返回的表格格式需要转换
                        table_data = []
                        for row in table:
                            if isinstance(row, list):
                                row_data = [str(cell) if cell else "" for cell in row]
                                table_data.append(row_data)
                        
                        if table_data:
                            if len(table_data) > 1:
                                df = pd.DataFrame(table_data[1:], columns=table_data[0])
                                tables.append({
                                    "data": df.to_dict(orient='records'),
                                    "columns": df.columns.tolist(),
                                    "table_index": table_idx,
                                })
                            else:
                                tables.append({
                                    "data": table_data,
                                    "columns": table_data[0] if table_data else [],
                                    "table_index": table_idx,
                                })
                    except Exception as e:
                        logger.warning(f"转换表格失败 (表{table_idx}): {e}")
                        tables.append({"data": str(table), "table_index": table_idx})
            
            return {
                "text": "\n\n".join(text_content),
                "tables": tables,
                "file_type": "doc",
            }
        except ImportError:
            logger.warning("docx2python未安装，DOC文件解析可能失败")
            # 备选方案：尝试使用antiword或textract（需要系统安装）
            raise ImportError("解析DOC文件需要安装docx2python: pip install docx2python")
    
    @staticmethod
    def _parse_pptx(file_path: str) -> Dict[str, Any]:
        """解析PPTX文件
        
        Args:
            file_path: PPTX文件路径
            
        Returns:
            解析后的内容
        """
        from pptx import Presentation
        
        prs = Presentation(file_path)
        text_content = []
        tables = []
        
        # 提取文本和表格
        for slide_num, slide in enumerate(prs.slides):
            slide_text = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_text.append(shape.text)
                
                # 提取表格
                if shape.shape_type == 19:  # MSO_SHAPE_TYPE.TABLE
                    table = shape.table
                    table_data = []
                    for row in table.rows:
                        row_data = [cell.text.strip() for cell in row.cells]
                        table_data.append(row_data)
                    
                    if table_data:
                        try:
                            if len(table_data) > 1:
                                df = pd.DataFrame(table_data[1:], columns=table_data[0])
                                tables.append({
                                    "data": df.to_dict(orient='records'),
                                    "columns": df.columns.tolist(),
                                    "slide": slide_num + 1,
                                })
                            else:
                                tables.append({
                                    "data": table_data,
                                    "columns": table_data[0] if table_data else [],
                                    "slide": slide_num + 1,
                                })
                        except Exception as e:
                            logger.warning(f"转换表格失败 (幻灯片{slide_num+1}): {e}")
                            tables.append({"data": table_data, "slide": slide_num + 1})
            
            if slide_text:
                text_content.append(f"\n幻灯片 {slide_num + 1}:\n" + "\n".join(slide_text))
        
        return {
            "text": "\n\n".join(text_content),
            "tables": tables,
            "slide_count": len(prs.slides),
            "file_type": "pptx",
        }
    
    @staticmethod
    def _parse_excel(file_path: str) -> Dict[str, Any]:
        """解析Excel文件
        
        Args:
            file_path: Excel文件路径
            
        Returns:
            解析后的内容
        """
        text_content = []
        tables = []
        
        try:
            # 读取所有工作表
            excel_file = pd.ExcelFile(file_path)
            
            for sheet_name in excel_file.sheet_names:
                df = pd.read_excel(excel_file, sheet_name=sheet_name)
                
                # 将每个工作表作为一个表格
                if not df.empty:
                    # 清理NaN值
                    df = df.fillna("")
                    
                    tables.append({
                        "data": df.to_dict(orient='records'),
                        "columns": df.columns.tolist(),
                        "sheet_name": sheet_name,
                    })
                    
                    # 同时生成文本描述
                    text_content.append(f"\n工作表: {sheet_name}\n")
                    text_content.append(df.to_string(index=False))
            
            return {
                "text": "\n\n".join(text_content),
                "tables": tables,
                "sheet_count": len(excel_file.sheet_names),
                "file_type": "excel",
            }
        except Exception as e:
            logger.error(f"解析Excel文件失败 {file_path}: {e}")
            raise
    
    @staticmethod
    def _parse_txt(file_path: str) -> Dict[str, Any]:
        """解析TXT文件（自动检测编码）
        
        Args:
            file_path: TXT文件路径
            
        Returns:
            解析后的内容
        """
        # 先检测编码
        with open(file_path, 'rb') as f:
            raw_data = f.read()
            result = chardet.detect(raw_data)
            encoding = result.get('encoding', 'utf-8')
            confidence = result.get('confidence', 0)
            
            logger.info(f"检测到编码: {encoding} (置信度: {confidence:.2f})")
        
        # 尝试多种编码
        encodings_to_try = [encoding, 'utf-8', 'gbk', 'gb2312', 'big5', 'latin-1']
        
        for enc in encodings_to_try:
            try:
                with open(file_path, 'r', encoding=enc) as f:
                    text = f.read()
                    return {
                        "text": text,
                        "tables": [],
                        "file_type": "txt",
                        "encoding": enc,
                    }
            except (UnicodeDecodeError, LookupError):
                continue
        
        # 如果所有编码都失败，使用错误处理模式
        logger.warning(f"无法确定文件编码，使用错误处理模式: {file_path}")
        with open(file_path, 'r', encoding='utf-8', errors='replace') as f:
            text = f.read()
            return {
                "text": text,
                "tables": [],
                "file_type": "txt",
                "encoding": "utf-8 (with errors replaced)",
            }



